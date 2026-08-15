"""Reconstruct PowerPoint slides as SVG, from geometry python-pptx already exposes.

Nothing rasterizes a real .pptx without either LibreOffice (a 400MB dependency, and a whole
extra container if it is to be memory-isolated), a proprietary library with a crippling free
tier, or a cloud API that would mean shipping tenant documents to a third party.

But this viewer is a citation-verification surface, not a presentation player. It has to show
"here is the slide, and here on it is the text that was cited". That needs faithful *geometry*,
not a faithful raster — and geometry is exactly what python-pptx gives away for free:
`slide_width`/`slide_height` and every shape's `left/top/width/height`, all in EMU
(914400 per inch). Placeholders inherit their position from the layout, and python-pptx
resolves that inheritance itself.

So each slide becomes an SVG: a background rect, an `<image>` per picture, and positioned
`<text>` per paragraph. The browser loads it through the same `<img>` tag as a PDF page image,
which is why the paged viewer needs no PPTX-specific code at all.

**What this does not reproduce:** charts, SmartArt, gradients, tables drawn as graphics,
WordArt, transitions, and exact font metrics. Text renders in a substituted font at the right
position and size. The UI labels the result as a reconstruction rather than pretending
otherwise — see the note in the viewer.

Every shape is rendered defensively. Theme-inherited fills raise on `.rgb` access, grouped
shapes nest arbitrarily, and any shape may have `None` geometry; one bad shape drops that
shape, never the slide, and never the ingestion.
"""

from __future__ import annotations

import base64
from typing import Any, Iterator
from xml.sax.saxutils import escape

# EMU (English Metric Units) per point. python-pptx reports every length in EMU; SVG user
# units here are points, so the whole file works in one coordinate system.
_EMU_PER_POINT = 12700

# Fallback slide box (10" x 7.5", the 4:3 default) for a deck whose dimensions are unreadable.
_DEFAULT_SLIDE_WIDTH_PT = 720.0
_DEFAULT_SLIDE_HEIGHT_PT = 540.0

_DEFAULT_FONT_SIZE_PT = 18.0
_FONT_STACK = "Segoe UI, Calibri, Helvetica, Arial, sans-serif"

# Pictures are inlined as data URIs so the SVG is a single self-contained object in storage.
# Past this size the base64 payload costs more than the fidelity is worth, and the picture is
# replaced by a neutral placeholder rect.
_MAX_INLINE_IMAGE_BYTES = 2 * 1024 * 1024


class SlideSvgRenderer:
    """Yields `(slide_number, svg_bytes, width_pt, height_pt)` for each slide."""

    shape = "paged"
    ext = "svg"
    mime = "image/svg+xml"

    def render(self, file_data: bytes) -> Iterator[tuple[int, bytes, float, float]]:
        import io

        from pptx import Presentation

        presentation = Presentation(io.BytesIO(file_data))
        width_pt = _to_points(presentation.slide_width, _DEFAULT_SLIDE_WIDTH_PT)
        height_pt = _to_points(presentation.slide_height, _DEFAULT_SLIDE_HEIGHT_PT)

        for index, slide in enumerate(presentation.slides, start=1):
            try:
                svg = self._slide_to_svg(slide, width_pt, height_pt)
            except Exception:  # noqa: BLE001 - a broken slide must not break the deck
                svg = _blank_slide(width_pt, height_pt)
            yield index, svg.encode("utf-8"), width_pt, height_pt

    # --- SVG assembly -------------------------------------------------------------

    def _slide_to_svg(self, slide: Any, width_pt: float, height_pt: float) -> str:
        body: list[str] = [
            f'<rect x="0" y="0" width="{width_pt:.2f}" height="{height_pt:.2f}" fill="#ffffff"/>'
        ]
        for shape in slide.shapes:
            body.extend(self._shape_to_svg(shape))

        return (
            f'<svg xmlns="http://www.w3.org/2000/svg" '
            f'xmlns:xlink="http://www.w3.org/1999/xlink" '
            f'viewBox="0 0 {width_pt:.2f} {height_pt:.2f}" '
            f'width="{width_pt:.2f}" height="{height_pt:.2f}">'
            f"{''.join(body)}"
            f"</svg>"
        )

    def _shape_to_svg(self, shape: Any) -> list[str]:
        try:
            # Grouped shapes nest; recurse so their children are placed too. Group children
            # already report absolute slide coordinates, so no transform is needed.
            if getattr(shape, "shape_type", None) is not None and hasattr(shape, "shapes"):
                parts: list[str] = []
                for child in shape.shapes:
                    parts.extend(self._shape_to_svg(child))
                return parts

            box = _shape_box(shape)
            if box is None:
                return []
            left, top, width, height = box

            if _is_picture(shape):
                return _picture_to_svg(shape, left, top, width, height)
            if getattr(shape, "has_text_frame", False):
                return _text_frame_to_svg(shape, left, top, width, height)
            return []
        except Exception:  # noqa: BLE001 - one unrenderable shape, not one unrenderable slide
            return []


# --- Helpers ----------------------------------------------------------------------


def _to_points(emu: Any, fallback: float) -> float:
    try:
        value = float(emu) / _EMU_PER_POINT
    except (TypeError, ValueError):
        return fallback
    return value if value > 0 else fallback


def _shape_box(shape: Any) -> tuple[float, float, float, float] | None:
    """Absolute position in points, or None when the shape has no usable geometry."""
    try:
        left, top = shape.left, shape.top
        width, height = shape.width, shape.height
    except Exception:  # noqa: BLE001
        return None
    if left is None or top is None or width is None or height is None:
        # A placeholder that inherits from a layout normally resolves via python-pptx; None
        # here means even the layout has no position, so there is nowhere to draw it.
        return None
    try:
        return (
            float(left) / _EMU_PER_POINT,
            float(top) / _EMU_PER_POINT,
            float(width) / _EMU_PER_POINT,
            float(height) / _EMU_PER_POINT,
        )
    except (TypeError, ValueError):
        return None


def _is_picture(shape: Any) -> bool:
    try:
        from pptx.enum.shapes import MSO_SHAPE_TYPE

        return shape.shape_type == MSO_SHAPE_TYPE.PICTURE
    except Exception:  # noqa: BLE001
        return False


def _picture_to_svg(
    shape: Any, left: float, top: float, width: float, height: float
) -> list[str]:
    try:
        image = shape.image
        blob = image.blob
        if not blob or len(blob) > _MAX_INLINE_IMAGE_BYTES:
            raise ValueError("image missing or too large to inline")
        content_type = image.content_type or "image/png"
        encoded = base64.b64encode(blob).decode("ascii")
        return [
            f'<image x="{left:.2f}" y="{top:.2f}" '
            f'width="{width:.2f}" height="{height:.2f}" '
            f'preserveAspectRatio="xMidYMid meet" '
            f'xlink:href="data:{content_type};base64,{encoded}"/>'
        ]
    except Exception:  # noqa: BLE001
        # A neutral placeholder keeps the slide's composition readable — the reader can still
        # see that something occupied this space, which is what matters for orientation.
        return [
            f'<rect x="{left:.2f}" y="{top:.2f}" '
            f'width="{width:.2f}" height="{height:.2f}" '
            f'fill="#f1f1f1" stroke="#d8d8d8" stroke-width="0.75"/>'
        ]


def _text_frame_to_svg(
    shape: Any, left: float, top: float, width: float, height: float
) -> list[str]:
    frame = shape.text_frame
    if frame is None:
        return []

    parts: list[str] = []
    # Start one line down: SVG `y` is the text baseline, so drawing at the box's top edge
    # would hang the glyphs above it.
    cursor = top
    for paragraph in frame.paragraphs:
        text = "".join(run.text or "" for run in paragraph.runs).strip()
        size_pt = _paragraph_size(paragraph)
        line_height = size_pt * 1.25
        cursor += line_height
        if not text:
            continue
        if cursor > top + height + line_height:
            # Overflowed its box — PowerPoint would have shrunk or clipped it. Stopping is
            # closer to the truth than spilling text across the rest of the slide.
            break
        anchor, x = _paragraph_anchor(paragraph, left, width)
        parts.append(
            f'<text x="{x:.2f}" y="{cursor:.2f}" '
            f'font-family="{_FONT_STACK}" font-size="{size_pt:.2f}" '
            f'text-anchor="{anchor}" fill="{_paragraph_color(paragraph)}" '
            f'{_paragraph_weight(paragraph)}>{escape(text)}</text>'
        )
    return parts


def _paragraph_size(paragraph: Any) -> float:
    try:
        for run in paragraph.runs:
            if run.font.size is not None:
                return max(6.0, float(run.font.size.pt))
        if paragraph.font.size is not None:
            return max(6.0, float(paragraph.font.size.pt))
    except Exception:  # noqa: BLE001
        pass
    return _DEFAULT_FONT_SIZE_PT


def _paragraph_weight(paragraph: Any) -> str:
    try:
        for run in paragraph.runs:
            if run.font.bold:
                return 'font-weight="600"'
    except Exception:  # noqa: BLE001
        pass
    return ""


def _paragraph_color(paragraph: Any) -> str:
    """The run colour when it is a literal RGB value, else near-black.

    Theme colours raise on `.rgb` — resolving them means walking the slide master's colour
    map, which is a lot of machinery for a reconstruction that is already approximate.
    """
    try:
        for run in paragraph.runs:
            color = run.font.color
            if color is not None and color.rgb is not None:
                return f"#{str(color.rgb)}"
    except Exception:  # noqa: BLE001
        pass
    return "#1a1a1a"


def _paragraph_anchor(paragraph: Any, left: float, width: float) -> tuple[str, float]:
    try:
        from pptx.enum.text import PP_ALIGN

        alignment = paragraph.alignment
        if alignment == PP_ALIGN.CENTER:
            return "middle", left + width / 2
        if alignment == PP_ALIGN.RIGHT:
            return "end", left + width
    except Exception:  # noqa: BLE001
        pass
    return "start", left


def _blank_slide(width_pt: float, height_pt: float) -> str:
    return (
        f'<svg xmlns="http://www.w3.org/2000/svg" '
        f'viewBox="0 0 {width_pt:.2f} {height_pt:.2f}" '
        f'width="{width_pt:.2f}" height="{height_pt:.2f}">'
        f'<rect x="0" y="0" width="{width_pt:.2f}" height="{height_pt:.2f}" fill="#ffffff"/>'
        f"</svg>"
    )
