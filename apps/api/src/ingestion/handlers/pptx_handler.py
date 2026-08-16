from __future__ import annotations

import io
from typing import Callable

from langchain_core.documents import Document

from src.core.text import sanitize_text_for_storage
from src.domain.entities import AssetStatus, KnowledgeAsset, RawContent
from src.domain.interfaces import IFileStorage

# Callable seam so tests can stub the deck reader without python-pptx or a real .pptx file,
# mirroring how YouTubeSourceHandler takes its transcript fetcher.
DeckReader = Callable[[bytes], list[dict]]


def _default_deck_reader(data: bytes) -> list[dict]:
    """Read a .pptx into `[{number, title, lines, notes}]`, one entry per slide.

    python-pptx is pure Python over lxml/Pillow — no compilation, no runtime, ~4MB on the
    image. Speaker notes are included because a deck's argument often lives there rather
    than on the slide.
    """
    from pptx import Presentation

    presentation = Presentation(io.BytesIO(data))
    slides: list[dict] = []

    for index, slide in enumerate(presentation.slides, start=1):
        lines: list[str] = []
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            for paragraph in shape.text_frame.paragraphs:
                text = "".join(run.text for run in paragraph.runs).strip()
                if text:
                    lines.append(text)

        notes = ""
        if slide.has_notes_slide:
            notes_frame = slide.notes_slide.notes_text_frame
            if notes_frame is not None:
                notes = (notes_frame.text or "").strip()

        slides.append(
            {
                "number": index,
                "title": lines[0] if lines else "",
                "lines": lines,
                "notes": notes,
            }
        )

    return slides


class PptxSourceHandler:
    """Source handler for uploaded PowerPoint decks (implements `ISourceHandler`).

    One `Document` per slide, with a `slide` locator — so "Slide 12" is what a citation says
    and what the viewer opens to. Slide titles are collected into `metadata["slides"]` so
    source detail can show the deck's structure without re-parsing the file.
    """

    def __init__(self, file_storage: IFileStorage, deck_reader: DeckReader = _default_deck_reader) -> None:
        self.file_storage = file_storage
        self.deck_reader = deck_reader

    def acquire(self, asset: KnowledgeAsset) -> RawContent:
        data = self.file_storage.download(asset.storage_key)
        return RawContent(
            data=data,
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        )

    def parse(self, asset: KnowledgeAsset, raw: RawContent) -> KnowledgeAsset:
        data = raw.data if isinstance(raw.data, bytes) else raw.data.encode("utf-8")

        try:
            slides = self.deck_reader(data)
        except ValueError:
            raise
        except Exception as exc:  # noqa: BLE001 - normalize any reader error to a clear message
            raise ValueError(f"This presentation could not be opened: {exc}") from exc

        if not slides:
            raise ValueError("This presentation has no slides in it")

        documents: list[Document] = []
        outline: list[dict] = []
        text_parts: list[str] = []

        for slide in slides:
            body = "\n".join(slide.get("lines") or [])
            notes = slide.get("notes") or ""
            if notes:
                body = f"{body}\n\nSpeaker notes: {notes}" if body else f"Speaker notes: {notes}"
            body = sanitize_text_for_storage(body).strip()
            number = slide.get("number")

            if body:
                documents.append(
                    Document(
                        page_content=body,
                        metadata={"locator": {"type": "slide", "value": number}},
                    )
                )
                text_parts.append(body)

            outline.append(
                {
                    "number": number,
                    "title": sanitize_text_for_storage(slide.get("title") or "").strip(),
                }
            )

        if not documents:
            raise ValueError("This presentation has no readable text on any slide")

        title = self._title(outline, asset.filename)

        return KnowledgeAsset(
            id=asset.id,
            knowledge_base_id=asset.knowledge_base_id,
            lineage_id=asset.lineage_id,
            version=asset.version,
            filename=asset.filename,
            title=title,
            source_type=asset.source_type,
            storage_key=asset.storage_key,
            status=AssetStatus.EXTRACTING,
            text_content="\n\n".join(text_parts),
            documents=documents,
            metadata={
                "filename": asset.filename,
                "title": title,
                "source_type": asset.source_type,
                "format": "slides",
                "content_type": raw.mime or asset.metadata.get("content_type"),
                "slide_count": len(slides),
                "slides": outline,
            },
        )

    def _title(self, outline: list[dict], filename: str) -> str:
        """The first slide's title is the deck's title; otherwise the filename stem."""
        for slide in outline:
            if slide.get("title"):
                return str(slide["title"])
        stem = filename.rsplit(".", 1)[0] if "." in filename else filename
        return sanitize_text_for_storage(stem or filename)
